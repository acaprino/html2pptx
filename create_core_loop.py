"""Create TIA Core Loop diagram as native PPTX shapes."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import math

# ── Presentation ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLD = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Colors ──
ORANGE   = RGBColor(0xFF, 0x62, 0x00)
GRAY     = RGBColor(0xD0, 0xD0, 0xD0)
LT_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x55, 0x55, 0x55)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)

# ── Dimensions (EMU) ──
# Angles: 0°=right(3 o'clock), positive=CW, matching OOXML and screen coords
CX       = prs.slide_width // 2
CY       = prs.slide_height // 2 + Inches(0.15)
OUTER_R  = Inches(1.65)
INNER_R  = Inches(1.22)       # white mask → ring thickness ≈ 0.43"
CENTER_R = Inches(0.78)
NODE_R   = Inches(0.26)

# ── Helpers ──

def _xy(angle_deg, r=OUTER_R):
    """Point on circle. 0°=right, CW positive (OOXML/screen convention)."""
    a = math.radians(angle_deg)
    return CX + int(r * math.cos(a)), CY + int(r * math.sin(a))


def _set_adj(shape, pairs):
    """Set shape geometry adjustment values via direct XML manipulation."""
    pg = shape._element.spPr.find(qn('a:prstGeom'))
    av = pg.find(qn('a:avLst'))
    for c in list(av):
        av.remove(c)
    for name, val in pairs:
        g = etree.SubElement(av, qn('a:gd'))
        g.set('name', name)
        g.set('fmla', f'val {int(val)}')


def add_shadow(shape, blur=76200, dist=25400, direction=5400000, alpha=15000):
    """Add subtle outer shadow to a shape."""
    spPr = shape._element.spPr
    # Remove existing effectLst if any
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None:
        spPr.remove(existing)
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', str(direction))
    outerShdw.set('algn', 'ctr')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', '000000')
    etree.SubElement(srgb, qn('a:alpha')).set('val', str(alpha))


def oval(cx, cy, r, fill=None, border=None, bw=Pt(0), shadow=False):
    s = SLD.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    if shadow:
        add_shadow(s)
    return s


def txt(cx, cy, w, h, text, sz, color, bold=False, fn='Segoe UI', anchor=MSO_ANCHOR.MIDDLE):
    t = SLD.shapes.add_textbox(cx - w // 2, cy - h // 2, w, h)
    tf = t.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Vertical centering
    txBody = t._element.find(qn('a:txBody'))
    if txBody is not None:
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = sz
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = fn
    return t


# ══════════════════════════════════════
#         BUILD TIA CORE LOOP
# ══════════════════════════════════════

# ── 1. RING: layered circles + PIE overlay ──
# Bottom layer: full orange circle
orange_base = oval(CX, CY, OUTER_R, fill=ORANGE)
add_shadow(orange_base, blur=100000, dist=30000, alpha=12000)

# Gray PIE sector overlay (top ~120°, from 10-o'clock to 2-o'clock through 12)
# OOXML PIE: adj1=start angle, adj2=end angle; CW from start to end
# 10 o'clock = 210°, 2 o'clock = 330° (in 0°=right, CW system)
pie_d = OUTER_R * 2
pie = SLD.shapes.add_shape(MSO_SHAPE.PIE, CX - OUTER_R, CY - OUTER_R, pie_d, pie_d)
_set_adj(pie, [('adj1', 210 * 60000), ('adj2', 330 * 60000)])
pie.fill.solid()
pie.fill.fore_color.rgb = GRAY
pie.line.fill.background()

# White inner circle (masks center → creates ring effect)
oval(CX, CY, INNER_R, fill=WHITE)

# ── 2. CENTER CIRCLE + TEXT ──
center = oval(CX, CY, CENTER_R, fill=WHITE, border=ORANGE, bw=Pt(2.5))

# "TIA" title
txt(CX, CY - Inches(0.15), Inches(1.4), Inches(0.5),
    'TIA', Pt(34), DARK, bold=True, fn='Poppins')

# "CORE LOOP" subtitle
txt(CX, CY + Inches(0.28), Inches(1.4), Inches(0.28),
    'CORE  LOOP', Pt(11), MID_GRAY, bold=False, fn='Poppins')

# ── 3. NODE ICONS (on the ring outer edge) ──
# Nodes divide the ring into 3 equal 120° segments

# 3a. TOP node (270° = 12 o'clock): "observe" — concentric target icon
tx, ty = _xy(270)
oval(tx, ty, NODE_R, fill=WHITE, border=MID_GRAY, bw=Pt(2.0))
oval(tx, ty, Inches(0.11), border=DARK, bw=Pt(1.5))
oval(tx, ty, Inches(0.04), fill=DARK)

# 3b. BOTTOM-LEFT node (150° ≈ 8 o'clock): "validate" — checkmark
bx, by = _xy(150)
oval(bx, by, NODE_R, fill=WHITE, border=ORANGE, bw=Pt(2.0))
txt(bx, by, Inches(0.4), Inches(0.4), '\u2713', Pt(18), ORANGE, bold=True)

# 3c. BOTTOM-RIGHT node (30° ≈ 4 o'clock): "execute" — lightning bolt
ex, ey = _xy(30)
oval(ex, ey, NODE_R, fill=WHITE, border=ORANGE, bw=Pt(2.0))
txt(ex, ey, Inches(0.4), Inches(0.4), '\u26A1', Pt(16), ORANGE, bold=True)

# ── 4. DIRECTION ARROWS (small triangles between nodes, pointing CW) ──
# Placed slightly outside the ring at the midpoint angles between nodes
# Triangle default = pointing UP; rotation = (angle + 180) % 360 for CW tangent
arrow_r = OUTER_R + Inches(0.18)
for ang in [330, 90, 210]:  # between nodes: 2 o'clock, 6 o'clock, 10 o'clock
    ax, ay = _xy(ang, arrow_r)
    sz = Inches(0.14)
    s = SLD.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                              ax - sz // 2, ay - sz // 2, sz, sz)
    s.fill.solid()
    s.fill.fore_color.rgb = GRAY
    s.line.fill.background()
    s.rotation = (ang + 180) % 360

# ── 5. WIFI / SIGNAL ARCS above top node ──
# Concentric arcs above the "observe" node, suggesting monitoring/awareness
for i, r in enumerate([Inches(0.30), Inches(0.42), Inches(0.55)]):
    arc_y = ty - Inches(0.20)  # shift arcs upward from node center
    d = r * 2
    s = SLD.shapes.add_shape(MSO_SHAPE.ARC, tx - r, arc_y - r, d, d)
    s.fill.background()
    c = LT_GRAY if i < 2 else RGBColor(0xEE, 0xEE, 0xEE)
    s.line.color.rgb = c
    s.line.width = Pt(2.0 - i * 0.4)
    # Show only top portion of arc (~220° to ~320°)
    _set_adj(s, [('adj1', 220 * 60000), ('adj2', 320 * 60000)])
    if i == 2:
        # Outermost arc: dashed
        s.line.dash_style = 3  # MSO_LINE_DASH_STYLE.DASH = 3

# ── Save ──
OUT = 'CoreLoop.pptx'
prs.save(OUT)
print(f'Saved: {OUT}')
