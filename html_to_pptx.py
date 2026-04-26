#!/usr/bin/env python3
"""
html_to_pptx.py - Parse HTML slides into native PPTX elements.

Renders each HTML in Playwright Chromium, extracts DOM element positions and
computed styles via injected JavaScript, then rebuilds the slide as native
python-pptx shapes, text boxes, and images. SVGs that can be expressed as
primitives (circle, rect, line, path, polygon, text) are converted natively;
complex SVGs fall back to a screenshot. Font Awesome glyphs are screenshotted.

The conversion is configured by SlideContext (no module-global mutation), so
this module is safe to import as a library: call convert(...) directly.
"""
from __future__ import annotations

import argparse
import logging
import os
import re
import sys
import tempfile
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any, Optional

try:
    from typing import TypedDict, NotRequired
except ImportError:  # pragma: no cover - Python < 3.11 fallback
    from typing_extensions import TypedDict, NotRequired  # type: ignore

from playwright.sync_api import (
    sync_playwright,
    Error as PlaywrightError,
    TimeoutError as PlaywrightTimeoutError,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE


# ── Logging ───────────────────────────────────────────────────────────
log = logging.getLogger("html2pptx")


# ── Named constants (all magic numbers consolidated) ──────────────────

# Default viewport (CSS pixels). Slide aspect ratio derives from this.
DEFAULT_VP_W = 1280
DEFAULT_VP_H = 720
STANDARD_SLIDE_H_INCHES = 7.5

# Unit conversion: CSS px -> typographic pt at 96 DPI.
PX_TO_PT = 72 / 96

# Slide layout index (PowerPoint default master ships index 6 as Blank).
BLANK_LAYOUT_INDEX = 6

# Element-size minimums (px) below which we drop silently.
MIN_SHAPE_PX = 5
MIN_TEXT_PX = 2
MIN_IMG_PX = 2
MIN_SVG_RECT_PX = 2
MIN_FREEFORM_BBOX_PX = 0.5
SVG_RECT_RX_THRESHOLD_PX = 2
ROUNDED_RECT_THRESHOLD_PX = 4
SVG_TEXT_PADDING_PX = 4

# Visibility thresholds.
MIN_VISIBLE_FILL_ALPHA = 0.15
MIN_SVG_ALPHA = 0.05
MIN_OPACITY_VISIBLE = 0.01
MIN_BORDER_PX = 1.0
MIN_CONTAINER_BG_ALPHA = 0.1

# Container detection / clipping geometry.
CONTAINER_MIN_AREA_PX = 100  # 10x10
BG_DEDUP_MARGIN_PX = 20
TEXT_OVERFLOW_DROP_RATIO = 0.5
SVG_FULL_SLIDE_THRESHOLD = 0.8

# Font handling.
DEFAULT_FONT_SIZE_PX = 16
DEFAULT_FONT_WEIGHT = 400
BOLD_WEIGHT_THRESHOLD = 600
FONT_SAFETY_MARGIN = 1.05
MIN_TEXT_HEIGHT_FACTOR = 1.2

# Wait timeouts (ms).
NETWORK_IDLE_TIMEOUT_MS = 3000
DEFAULT_TAILWIND_WAIT_MS = 1500
DEFAULT_TAILWIND_MIN_MS = 50

# SVG primitive sampling.
SVG_MAX_PRIMITIVES = 500
SVG_PATH_MIN_PTS = 4
SVG_PATH_MAX_PTS = 200
SVG_PATH_LENGTH_DIVISOR = 3
PATH_MIN_LENGTH_PX = 0.5
POLYGON_MAX_POINTS = 200

# Stacking-trick depth offset for rectangular borders (paints between parent and child).
BORDER_DEPTH_OFFSET = 0.5

# Network egress allowlist (default; configurable via --allow-network).
DEFAULT_NETWORK_ALLOWLIST = [
    "cdn.tailwindcss.com",
    "fonts.googleapis.com",
    "fonts.gstatic.com",
]
# file:// URLs are always allowed (the slide itself).

# Input hardening.
MAX_INPUT_HTML_BYTES = 5 * 1024 * 1024  # 5 MB per slide
MAX_SLIDES_DEFAULT = 5000

# Tailwind v3 patch CDN (pinned). Loaded with SRI when available.
TAILWIND_CDN_URL = "https://cdn.tailwindcss.com/3.4.17?plugins=forms,typography"
TAILWIND_CDN_TAG = f'<script src="{TAILWIND_CDN_URL}"></script>'
# Stricter than legacy: anchor on tailwindcss@2 path; bound character classes to
# prevent catastrophic backtracking on hostile near-matches.
TAILWIND_V2_LINK_RE = re.compile(
    r'<link\s[^>]{0,500}?href=["\']?[^"\'>]{0,400}?tailwindcss@2[^"\'>]{0,400}?["\'][^>]{0,200}?>',
    re.IGNORECASE,
)

# Flex overflow workaround.
FLEX_FIX_CSS = "<style>.flex-1{min-height:0!important;min-width:0!important;}</style>"

# Web font -> width compensation ratio (web vs Windows fallback) and substitute font.
FONT_RATIOS_HARDCODED = {
    "Poppins": 1.137,
    "Inter": 1.08,
    "Roboto Mono": 1.092,
}
FONT_MAP = {
    "Poppins": "Segoe UI",
    "Inter": "Segoe UI",
    "Roboto Mono": "Consolas",
}

ALIGN = {
    "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
    "left": PP_ALIGN.LEFT, "start": PP_ALIGN.LEFT,
    "end": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
}

ROOT = Path(__file__).resolve().parent
DEFAULT_HTML_DIR = ROOT / "presentazione_html"
DEFAULT_OUTPUT = ROOT / "Slides1.pptx"


# ── CSS simplification levels (cumulative) ────────────────────────────
SIMPLIFY_LEVELS = [
    (1, "*, *::before, *::after { animation: none !important; transition: none !important; "
        "animation-duration: 0s !important; }"),
    (3, "* { box-shadow: none !important; text-shadow: none !important; }"),
    (5, "* { filter: none !important; backdrop-filter: none !important; "
        "-webkit-backdrop-filter: none !important; }"),
    (7, "* { background-image: none !important; }"),
    (9, "* { outline: none !important; } "
        "*::before, *::after { content: none !important; }"),
]


def _simplify_css(level: int) -> str:
    """Generate <style> tag with CSS overrides for the given simplification level."""
    if level <= 0:
        return ""
    rules = [css for threshold, css in SIMPLIFY_LEVELS if level >= threshold]
    return "<style>" + "\n".join(rules) + "</style>"


# ── Walker schema (TypedDicts; documents JS<->Python contract) ────────

class RunDict(TypedDict, total=False):
    t: str
    ff: str
    fs: float
    fw: int
    fi: bool
    co: str
    tt: str
    hlBg: NotRequired[str]
    hlX: NotRequired[float]
    hlY: NotRequired[float]
    hlW: NotRequired[float]
    hlH: NotRequired[float]
    hlBr: NotRequired[float]


class ShapeEntry(TypedDict, total=False):
    x: float
    y: float
    w: float
    h: float
    dp: float
    seq: int
    z: int                       # CSS z-index (0 if static)
    bg: str
    op: float
    br: float                    # legacy single radius (kept for compatibility)
    brX: NotRequired[float]      # per-axis horizontal radius (CSS-spec)
    brY: NotRequired[float]      # per-axis vertical radius (CSS-spec)
    circ: NotRequired[bool]
    bco: NotRequired[str]
    bw: NotRequired[float]


class TextEntry(TypedDict, total=False):
    x: float
    y: float
    w: float
    h: float
    tx: float
    tw: float
    dp: float
    seq: int
    z: int
    runs: list[RunDict]
    ta: str
    multiline: bool              # browser-precise: getClientRects().length > 1


class ScreenshotEntry(TypedDict, total=False):
    x: float
    y: float
    w: float
    h: float
    dp: float
    seq: int
    z: int
    i: int
    cw: float                    # container width at extraction time (used for skip threshold)
    ch: float


class SvgPrimitive(TypedDict, total=False):
    type: str
    x: float
    y: float
    w: float
    h: float
    dp: float
    seq: int
    fill: Optional[str]
    stroke: Optional[str]
    strokeWidth: float
    dashed: bool
    opacity: float
    fillOpacity: float
    strokeOpacity: float
    rx: NotRequired[float]
    points: NotRequired[list[list[float]]]
    closed: NotRequired[bool]
    text: NotRequired[str]
    fontFamily: NotRequired[str]
    fontSize: NotRequired[float]
    fontWeight: NotRequired[int]
    textAnchor: NotRequired[str]


class WalkerOutput(TypedDict, total=False):
    bg: str
    bodyBg: str
    cw: float
    ch: float
    shapes: list[ShapeEntry]
    texts: list[TextEntry]
    svgs: list[ScreenshotEntry]
    icons: list[ScreenshotEntry]
    nativeSvgs: list[dict[str, Any]]
    fontRatios: dict[str, dict[str, Any]]


# ── SlideContext: replaces the legacy mutable module globals ──────────

@dataclass
class SlideContext:
    """All per-conversion configuration.

    Replaces the legacy module globals (VP_W, VP_H, SLIDE_W, SLIDE_H, SCALE)
    that used to be mutated inside main(). Helpers take ``ctx`` explicitly.
    """
    vp_w: int
    vp_h: int
    slide_w: int   # EMU
    slide_h: int   # EMU
    scale: float   # EMU per CSS pixel

    @classmethod
    def from_viewport(cls, vp_w: int, vp_h: int) -> "SlideContext":
        if vp_w <= 0 or vp_h <= 0:
            raise ValueError(f"viewport must be positive, got {vp_w}x{vp_h}")
        slide_h = Inches(STANDARD_SLIDE_H_INCHES)
        # Use exact integer arithmetic so 16:9 lands on PowerPoint's canonical
        # 12,192,000 EMU rather than Inches(13.333) -> 12,191,904.
        slide_w = int(round(slide_h * vp_w / vp_h))
        scale = slide_w / vp_w
        return cls(vp_w=vp_w, vp_h=vp_h, slide_w=slide_w, slide_h=slide_h, scale=scale)

    def px(self, v: float) -> int:
        """CSS pixels to EMU. Round (not truncate) so adjacent elements align."""
        return int(round(v * self.scale))

    def pt(self, v: float) -> Pt:
        """CSS pixels to typographic points (96 DPI assumption)."""
        return Pt(v * PX_TO_PT)


# ── ConversionReport: aggregated failure tracking ─────────────────────

@dataclass
class ConversionReport:
    slides_total: int = 0
    slides_processed: int = 0
    slides_skipped: int = 0
    elements_skipped: int = 0
    screenshots_failed: int = 0
    networkidle_timeouts: int = 0
    tailwind_wait_timeouts: int = 0
    preprocess_failures: int = 0
    save_partial: bool = False
    failures_by_slide: dict[str, list[str]] = field(default_factory=dict)

    def warn(self, slide_name: str, message: str) -> None:
        log.warning("[%s] %s", slide_name, message)
        self.failures_by_slide.setdefault(slide_name, []).append(message)

    def has_failures(self) -> bool:
        return bool(self.failures_by_slide) or self.save_partial

    def summary_line(self) -> str:
        parts = [f"{self.slides_processed}/{self.slides_total} slides"]
        if self.slides_skipped:
            parts.append(f"{self.slides_skipped} skipped")
        if self.elements_skipped:
            parts.append(f"{self.elements_skipped} elements dropped")
        if self.screenshots_failed:
            parts.append(f"{self.screenshots_failed} screenshots failed")
        if self.networkidle_timeouts:
            parts.append(f"{self.networkidle_timeouts} networkidle timeouts")
        if self.tailwind_wait_timeouts:
            parts.append(f"{self.tailwind_wait_timeouts} Tailwind-wait timeouts")
        if self.preprocess_failures:
            parts.append(f"{self.preprocess_failures} preprocess failures")
        if self.save_partial:
            parts.append("PARTIAL save")
        return ", ".join(parts)


# ── Color parsing (CSS L4-aware, ReDoS-safe) ──────────────────────────

# Bounded digit ranges to prevent catastrophic backtracking. 5 leading digits
# accepts every valid CSS number (RGB channels are 0-255, hue is 0-360, etc.)
# while keeping the per-token match short enough to preclude pathological backtracking.
_RE_NUM = r"-?\d{1,5}(?:\.\d{1,5})?"
_RE_PCT = r"-?\d{1,5}(?:\.\d{1,5})?%"
_RE_NUM_OR_PCT = rf"(?:{_RE_PCT}|{_RE_NUM})"
_RE_ALPHA = rf"(?:{_RE_NUM}%?)"

_RGB_RE = re.compile(
    rf"^rgba?\(\s*({_RE_NUM_OR_PCT})\s*[,\s]\s*({_RE_NUM_OR_PCT})\s*[,\s]\s*({_RE_NUM_OR_PCT})"
    rf"(?:\s*[,/]\s*({_RE_ALPHA}))?\s*\)$",
    re.IGNORECASE,
)
_HSL_RE = re.compile(
    rf"^hsla?\(\s*({_RE_NUM})(?:deg)?\s*[,\s]\s*({_RE_PCT})\s*[,\s]\s*({_RE_PCT})"
    rf"(?:\s*[,/]\s*({_RE_ALPHA}))?\s*\)$",
    re.IGNORECASE,
)
_HEX_RE = re.compile(r"^#([0-9a-fA-F]{3,8})$")

_NAMED_COLORS: dict[str, tuple[int, int, int, float]] = {
    "transparent": (0, 0, 0, 0.0),
    "white": (255, 255, 255, 1.0),
    "black": (0, 0, 0, 1.0),
    "red": (255, 0, 0, 1.0),
    "green": (0, 128, 0, 1.0),
    "blue": (0, 0, 255, 1.0),
    "gray": (128, 128, 128, 1.0),
    "grey": (128, 128, 128, 1.0),
}


def _parse_channel(raw: str) -> float:
    raw = raw.strip()
    if raw.endswith("%"):
        return float(raw[:-1]) * 255.0 / 100.0
    return float(raw)


def _parse_alpha(raw: Optional[str]) -> float:
    if raw is None:
        return 1.0
    raw = raw.strip()
    if raw.endswith("%"):
        return float(raw[:-1]) / 100.0
    return float(raw)


def _hsl_to_rgb(h: float, s_pct: float, l_pct: float) -> tuple[float, float, float]:
    h = (h % 360) / 360.0
    s = max(0.0, min(1.0, s_pct / 100.0))
    lum = max(0.0, min(1.0, l_pct / 100.0))
    if s == 0:
        v = lum * 255.0
        return v, v, v
    q = lum * (1 + s) if lum < 0.5 else lum + s - lum * s
    p = 2 * lum - q

    def hue(t: float) -> float:
        t = t % 1.0
        if t < 1 / 6:
            return p + (q - p) * 6 * t
        if t < 1 / 2:
            return q
        if t < 2 / 3:
            return p + (q - p) * (2 / 3 - t) * 6
        return p

    return (hue(h + 1 / 3) * 255.0, hue(h) * 255.0, hue(h - 1 / 3) * 255.0)


def _clamp_rgba(r: float, g: float, b: float, a: float) -> tuple[int, int, int, float]:
    return (
        max(0, min(255, int(round(r)))),
        max(0, min(255, int(round(g)))),
        max(0, min(255, int(round(b)))),
        max(0.0, min(1.0, a)),
    )


def parse_rgba(raw: Optional[str]) -> Optional[tuple[int, int, int, float]]:
    """Parse a CSS color string to (R, G, B, A).

    Supports rgb()/rgba() (incl. CSS L4 space-separated), hsl()/hsla(),
    hex (#rgb, #rgba, #rrggbb, #rrggbbaa), and a small set of named colors.
    Logs and returns None for CSS L4 functions we cannot accurately convert
    (oklch, oklab, color(), color-mix, lab, lch) so the caller's silent-skip
    paths are at least observable.
    """
    if not raw:
        return None
    s = raw.strip()
    if not s:
        return None

    if s in _NAMED_COLORS:
        r, g, b, a = _NAMED_COLORS[s]
        return (int(r), int(g), int(b), a)

    m = _RGB_RE.match(s)
    if m:
        r = _parse_channel(m[1])
        g = _parse_channel(m[2])
        b = _parse_channel(m[3])
        a = _parse_alpha(m[4])
        return _clamp_rgba(r, g, b, a)

    m = _HSL_RE.match(s)
    if m:
        h = float(m[1])
        s_p = float(m[2].rstrip("%"))
        l_p = float(m[3].rstrip("%"))
        a = _parse_alpha(m[4])
        r, g, b = _hsl_to_rgb(h, s_p, l_p)
        return _clamp_rgba(r, g, b, a)

    m = _HEX_RE.match(s)
    if m:
        h = m[1]
        if len(h) in (3, 4):
            h = "".join(c * 2 for c in h)
        if len(h) == 6:
            return _clamp_rgba(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16), 1.0)
        if len(h) == 8:
            return _clamp_rgba(
                int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16),
                int(h[6:8], 16) / 255.0,
            )

    if s.startswith(("oklch(", "oklab(", "color(", "color-mix(", "lab(", "lch(")):
        log.debug("parse_rgba: unsupported CSS L4 color: %s", s)
    return None


def to_rgb(c: Optional[tuple[int, int, int, float]]) -> Optional[RGBColor]:
    if not c:
        return None
    return RGBColor(c[0], c[1], c[2])


def _strip_dangerous_unicode(s: str) -> str:
    """Strip bidi-override and zero-width characters (defense-in-depth)."""
    return re.sub(r"[‪-‮⁦-⁩​-‏]", "", s)


def _apply_text_transform(text: str, tt: str) -> str:
    """Apply CSS text-transform.

    CSS-spec ``capitalize`` capitalizes first letter of each word and leaves
    the rest untouched -- different from Python ``str.capitalize`` which
    lowercases all but the first character of the string.
    """
    if tt == "uppercase":
        return text.upper()
    if tt == "lowercase":
        return text.lower()
    if tt == "capitalize":
        return " ".join(w[:1].upper() + w[1:] for w in text.split(" "))
    return text


# ── Color/border helpers (consolidate the parse_rgba->to_rgb idiom) ──

def _apply_solid_fill(shape, color_str: Optional[str], op: float = 1.0) -> bool:
    c = parse_rgba(color_str)
    if not c or (op * c[3]) < MIN_VISIBLE_FILL_ALPHA:
        return False
    shape.fill.solid()
    shape.fill.fore_color.rgb = to_rgb(c)
    return True


def _apply_line_color(shape, color_str: Optional[str], width_px: float, dashed: bool = False) -> bool:
    c = parse_rgba(color_str)
    if not c:
        shape.line.fill.background()
        return False
    shape.line.color.rgb = to_rgb(c)
    shape.line.width = Pt(max(width_px * PX_TO_PT, 0.25))
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return True


def _resolve_font(family: str, font_ratios: dict) -> tuple[str, float]:
    """Return (substitute_font_name, width_compensation_factor).

    Browser-measured ratio first, hardcoded fallback second, no compensation last.
    The 5% safety margin (FONT_SAFETY_MARGIN) is applied uniformly.
    """
    measured = font_ratios.get(family, {})
    fallback_name = measured.get("fb") or FONT_MAP.get(family, family)
    ratio = measured.get("r") or FONT_RATIOS_HARDCODED.get(family, 1.0)
    factor = (1.0 / ratio) * FONT_SAFETY_MARGIN
    return fallback_name, factor


# ── Path validation (input/output hardening) ──────────────────────────

def _safe_input_dir(raw: str) -> Path:
    p = Path(raw).resolve(strict=False)
    if not p.exists():
        raise SystemExit(f"--input not found: {p}")
    if not p.is_dir():
        raise SystemExit(f"--input is not a directory: {p}")
    if p.is_symlink():
        raise SystemExit(f"--input must not be a symlink: {p}")
    return p


def _safe_output_path(raw: str) -> Path:
    p = Path(raw)
    s = str(p)
    if s.startswith("\\\\") or s.startswith("//"):
        raise SystemExit("--output must not be a UNC path")
    if ":" in p.name:
        raise SystemExit("--output filename must not contain ':'")
    if p.suffix.lower() != ".pptx":
        raise SystemExit("--output must have .pptx extension")
    parent = p.parent.resolve(strict=False) if p.parent != Path("") else Path.cwd()
    if not parent.exists():
        raise SystemExit(f"--output parent directory does not exist: {parent}")
    return p


def _safe_input_file(p: Path, html_dir: Path) -> bool:
    """Verify p resolves inside html_dir (no symlink escape) and is small enough."""
    try:
        rp = p.resolve(strict=True)
    except OSError:
        return False
    try:
        rp.relative_to(html_dir.resolve(strict=False))
    except ValueError:
        return False
    if p.is_symlink():
        return False
    try:
        if p.stat().st_size > MAX_INPUT_HTML_BYTES:
            log.warning("[%s] file exceeds %d bytes; skipping", p.name, MAX_INPUT_HTML_BYTES)
            return False
    except OSError:
        return False
    return True


# ── Shared container detection JS (used by preprocess + extractor) ────

# Filters off-screen overlays (top: -9999px), opacity: 0, pointer-events: none,
# inert/aria-hidden -- in addition to display:none / visibility:hidden.
FIND_CONTAINER_JS = r"""
    var best = null, bestArea = 0;
    var ch = document.body.children;
    var vw = window.innerWidth, vh = window.innerHeight;
    for (var fci = 0; fci < ch.length; fci++) {
        var bodyChild = ch[fci], tag = bodyChild.tagName;
        if (!tag || tag === 'SCRIPT' || tag === 'STYLE' || tag === 'LINK' || tag === 'META') continue;
        if (bodyChild.hasAttribute && (bodyChild.hasAttribute('inert') || bodyChild.getAttribute('aria-hidden') === 'true')) continue;
        var s = getComputedStyle(bodyChild);
        if (s.display === 'none' || s.visibility === 'hidden') continue;
        if (parseFloat(s.opacity) < 0.01) continue;
        var r = bodyChild.getBoundingClientRect();
        // Reject elements wholly off-screen.
        if (r.right <= 0 || r.bottom <= 0 || r.left >= vw || r.top >= vh) continue;
        if (r.width < 1 || r.height < 1) continue;
        var area = r.width * r.height;
        if (area > bestArea) { bestArea = area; best = bodyChild; }
    }
"""

# Strip pre-existing data-si / data-fi from user HTML so the walker's tagging
# cannot be hijacked by attacker-supplied attributes (HIGH-14).
PREPROCESS_JS = r"""() => {
""" + FIND_CONTAINER_JS + r"""
    var c = bestArea >= 100 ? best : null;
    if (!c) return 0;

    // Strip pre-existing data-si / data-fi attributes from user HTML.
    var existing = c.querySelectorAll('[data-si],[data-fi]');
    for (var ex = 0; ex < existing.length; ex++) {
        existing[ex].removeAttribute('data-si');
        existing[ex].removeAttribute('data-fi');
    }

    // Resolve <br> into block-level wrappers so the walker emits a separate
    // text entry per visual line.
    var brCount = 0;
    var candidates = c.querySelectorAll('*');
    for (var ci = 0; ci < candidates.length; ci++) {
        var candidate = candidates[ci];
        var hasBr = false;
        for (var j = 0; j < candidate.childNodes.length; j++) {
            if (candidate.childNodes[j].nodeType === 1 && candidate.childNodes[j].tagName === 'BR') {
                hasBr = true; break;
            }
        }
        if (!hasBr) continue;

        var savedTextAlign = getComputedStyle(candidate).textAlign;
        var nodes = Array.from(candidate.childNodes);
        var segments = [[]];
        for (var ni = 0; ni < nodes.length; ni++) {
            var node = nodes[ni];
            if (node.nodeType === 1 && node.tagName === 'BR') {
                segments.push([]);
                brCount++;
            } else {
                segments[segments.length - 1].push(node);
            }
        }

        while (candidate.firstChild) candidate.removeChild(candidate.firstChild);
        for (var si = 0; si < segments.length; si++) {
            var seg = segments[si];
            var hasContent = false;
            for (var k = 0; k < seg.length; k++) {
                if (seg[k].nodeType === 3 && seg[k].textContent.trim()) { hasContent = true; break; }
                if (seg[k].nodeType === 1) { hasContent = true; break; }
            }
            if (!hasContent) continue;

            var wrapper = document.createElement('div');
            wrapper.style.textAlign = savedTextAlign;
            for (var m = 0; m < seg.length; m++) {
                wrapper.appendChild(seg[m]);
            }
            candidate.appendChild(wrapper);
        }
    }

    return brCount;
}"""


# ── Hide/restore for clean per-element screenshots (private symbol-keyed state).
# Avoids window.__ssHidden being clobbered by hostile HTML.
HIDE_RESTORE_INIT_JS = r"""
(() => {
    if (window.__h2pH) return;
    var stash = null;
    window.__h2pH = function(selector) {
        var target = document.querySelector(selector);
        if (!target) return 0;
        var hidden = [];
        var savedVis = [];
        var current = target;
        while (current && current.parentElement) {
            var parent = current.parentElement;
            for (var i = 0; i < parent.children.length; i++) {
                var sib = parent.children[i];
                if (sib !== current) {
                    savedVis.push(sib.style.visibility);
                    hidden.push(sib);
                    sib.style.visibility = 'hidden';
                }
            }
            current = parent;
            if (current === document.documentElement) break;
        }
        stash = { hidden: hidden, savedVis: savedVis };
        return hidden.length;
    };
    window.__h2pR = function() {
        if (!stash) return;
        for (var i = 0; i < stash.hidden.length; i++) {
            stash.hidden[i].style.visibility = stash.savedVis[i] || '';
        }
        stash = null;
    };
})();
"""

# Tailwind JIT readiness probe. Resolves quickly when CDN finishes; falls back
# to a brief timeout cap rather than the legacy unconditional 500 ms wait.
TAILWIND_READY_JS = r"""
() => {
    // Tailwind Play CDN injects a <style> with attribute "data-jit-style" or
    // sets window.tailwind once the runtime is ready. document.fonts.ready
    // gates web-font availability.
    return Promise.race([
        document.fonts.ready,
        new Promise(function(resolve) { setTimeout(resolve, 1); })
    ]).then(function() {
        return !!(window.tailwind || document.querySelector('style[data-tailwind], style[data-jit-style]'));
    });
}
"""


# ── DOM extraction (executed inside Playwright via page.evaluate) ─────

# Major changes from legacy:
#  - viewport-aware container detection
#  - per-axis border-radius (brX, brY) plus legacy br
#  - CSS z-index captured for each element
#  - text multiline flag from getClientRects().length (precise)
#  - container width/height emitted on each svg/icon for python-side threshold consistency
#  - inline-block / inline-flex / inline-grid treated as inline
#  - SVG path closed-flag handles multi-subpath (returns null to fall back)

EXTRACT_JS = r"""() => {
""" + FIND_CONTAINER_JS + r"""
    var c = bestArea >= 100 ? best : null;
    if (!c) return null;
    var cr = c.getBoundingClientRect();
    var ox = cr.left, oy = cr.top;
    var out = {
        bg: getComputedStyle(c).backgroundColor,
        bodyBg: getComputedStyle(document.body).backgroundColor,
        cw: cr.width, ch: cr.height,
        shapes: [], texts: [], svgs: [], icons: [], nativeSvgs: []
    };

    // Tag SVGs and FA icons with our own attributes (existing data-si/data-fi
    // were stripped in PREPROCESS_JS so we can rely on uniqueness).
    c.querySelectorAll('svg').forEach(function(el, i) { el.setAttribute('data-si', String(i)); });
    c.querySelectorAll('.fas,.far,.fab,.fa-solid,.fa').forEach(function(el, i) { el.setAttribute('data-fi', String(i)); });

    function gr(el) {
        var b = el.getBoundingClientRect();
        return { x: b.left - ox, y: b.top - oy, w: b.width, h: b.height };
    }
    function vis(el) {
        var b = el.getBoundingClientRect();
        if (b.width < 0.5 || b.height < 0.5) return false;
        var s = getComputedStyle(el);
        return s.display !== 'none' && s.visibility !== 'hidden' && parseFloat(s.opacity) > 0.01;
    }
    function inl(el) {
        var d = getComputedStyle(el).display;
        return d === 'inline' || d === 'inline-block' || d === 'inline-flex' || d === 'inline-grid';
    }
    function isFA(el) {
        if (!el.classList) return false;
        return ['fas','far','fab','fa-solid','fa'].some(function(cn) { return el.classList.contains(cn); });
    }
    function zIndexOf(s) {
        if (s.position === 'static' || s.zIndex === 'auto') return 0;
        var v = parseInt(s.zIndex, 10);
        return isNaN(v) ? 0 : v;
    }
    function sty(el) {
        var s = getComputedStyle(el);
        return {
            ff: s.fontFamily.split(',')[0].replace(/['"]/g, '').trim(),
            fs: parseFloat(s.fontSize),
            fw: parseInt(s.fontWeight) || 400,
            fi: s.fontStyle === 'italic',
            co: s.color,
            tt: s.textTransform
        };
    }

    function svgStyle(child) {
        var cs = getComputedStyle(child);
        var fill = cs.fill || '';
        var stroke = cs.stroke || '';
        var sw = parseFloat(cs.strokeWidth) || 0;
        var da = cs.strokeDasharray || '';
        var op = parseFloat(cs.opacity); if (isNaN(op)) op = 1;
        var fop = parseFloat(cs.fillOpacity); if (isNaN(fop)) fop = 1;
        var sop = parseFloat(cs.strokeOpacity); if (isNaN(sop)) sop = 1;
        var hasFill = fill && fill !== 'none' && fill.indexOf('url(') < 0;
        var hasStroke = stroke && stroke !== 'none' && sw > 0 && stroke.indexOf('url(') < 0;
        return {
            fill: hasFill ? fill : null,
            stroke: hasStroke ? stroke : null,
            strokeWidth: hasStroke ? sw : 0,
            dashed: da && da !== 'none',
            opacity: op, fillOpacity: fop, strokeOpacity: sop
        };
    }

    function transformPt(ctm, px, py) {
        return [ctm.a * px + ctm.c * py + ctm.e - ox,
                ctm.b * px + ctm.d * py + ctm.f - oy];
    }

    function extractSvgNative(svgEl, baseDp, baseSeq) {
        var sr = svgEl.getBoundingClientRect();
        if (sr.width > cr.width * 0.8 && sr.height > cr.height * 0.8) return null;
        var blockers = svgEl.querySelectorAll('pattern, clipPath, mask, filter, foreignObject, image, use');
        if (blockers.length > 0) return null;

        var prims = svgEl.querySelectorAll('circle, ellipse, rect, line, path, polygon, polyline, text');
        var elements = [];

        for (var i = 0; i < prims.length; i++) {
            var child = prims[i];
            if (child.closest('defs') || child.closest('marker')) continue;
            var ccs = getComputedStyle(child);
            if (ccs.display === 'none' || ccs.visibility === 'hidden') continue;
            if (parseFloat(ccs.opacity) < 0.01) continue;

            var tag = child.tagName.toLowerCase();
            var style = svgStyle(child);
            if (!style.fill && !style.stroke) continue;
            var br = child.getBoundingClientRect();
            if (br.width < 0.5 && br.height < 0.5) continue;
            if (elements.length >= 500) break;
            var base = {
                type: tag,
                x: br.left - ox, y: br.top - oy,
                w: br.width, h: br.height,
                dp: baseDp, seq: baseSeq + elements.length
            };
            base.fill = style.fill; base.stroke = style.stroke;
            base.strokeWidth = style.strokeWidth; base.dashed = style.dashed;
            base.opacity = style.opacity; base.fillOpacity = style.fillOpacity;
            base.strokeOpacity = style.strokeOpacity;

            if (tag === 'circle' || tag === 'ellipse' || tag === 'rect') {
                if (tag === 'rect') {
                    base.rx = parseFloat(child.getAttribute('rx')) || 0;
                }
                elements.push(base);

            } else if (tag === 'line') {
                var ctm = child.getScreenCTM(); if (!ctm) return null;
                var x1 = child.x1.baseVal.value, y1 = child.y1.baseVal.value;
                var x2 = child.x2.baseVal.value, y2 = child.y2.baseVal.value;
                base.points = [transformPt(ctm, x1, y1), transformPt(ctm, x2, y2)];
                base.closed = false;
                elements.push(base);

            } else if (tag === 'path') {
                try {
                    var d = child.getAttribute('d') || '';
                    // Multi-subpath paths (multiple M commands) cannot be expressed as
                    // a single PPTX freeform with one `closed` flag; fall back.
                    var dt = d.trim();
                    var moves = dt.match(/[Mm]/g);
                    if (moves && moves.length > 1) return null;
                    var totalLen = child.getTotalLength();
                    if (totalLen < 0.5) continue;
                    var ctm2 = child.getScreenCTM(); if (!ctm2) return null;
                    var nPts = Math.max(4, Math.min(200, Math.round(totalLen / 3)));
                    var pts = [];
                    for (var j = 0; j <= nPts; j++) {
                        var pt = child.getPointAtLength(totalLen * j / nPts);
                        pts.push(transformPt(ctm2, pt.x, pt.y));
                    }
                    base.points = pts;
                    base.closed = /[zZ]\s*$/.test(dt);
                    elements.push(base);
                } catch (e) { return null; }

            } else if (tag === 'polygon' || tag === 'polyline') {
                var ctm3 = child.getScreenCTM(); if (!ctm3) return null;
                var plist = child.points;
                if (!plist || plist.numberOfItems < 2) continue;
                if (plist.numberOfItems > 200) return null;
                var ppts = [];
                for (var k = 0; k < plist.numberOfItems; k++) {
                    var pp = plist.getItem(k);
                    ppts.push(transformPt(ctm3, pp.x, pp.y));
                }
                base.points = ppts;
                base.closed = (tag === 'polygon');
                elements.push(base);

            } else if (tag === 'text') {
                var ff = ccs.fontFamily || '';
                if (ff.toLowerCase().indexOf('fontawesome') >= 0 ||
                    ff.toLowerCase().indexOf('font awesome') >= 0) continue;
                base.text = child.textContent || '';
                base.fontFamily = ff.split(',')[0].replace(/['"]/g, '').trim();
                base.fontSize = parseFloat(ccs.fontSize) || 16;
                base.fontWeight = parseInt(ccs.fontWeight) || 400;
                base.textAnchor = child.getAttribute('text-anchor') || ccs.textAnchor || 'start';
                elements.push(base);
            }
        }
        if (elements.length === 0) return null;
        return { elements: elements, dp: baseDp, seq: baseSeq };
    }

    var seq = 0;
    var dp = 0;

    function walk(el) {
        var tag = el.tagName;
        if (!tag || ['SCRIPT','STYLE','LINK','META','HEAD','BR','HR'].indexOf(tag) >= 0) return;
        if (!vis(el)) return;
        var rect = gr(el);
        var cs = getComputedStyle(el);
        var z = zIndexOf(cs);

        if (tag === 'svg' || tag === 'SVG') {
            var si = el.getAttribute('data-si');
            if (si !== null) {
                var native = extractSvgNative(el, dp, seq);
                if (native) {
                    out.nativeSvgs.push(native);
                    seq += native.elements.length;
                } else {
                    out.svgs.push({ x: rect.x, y: rect.y, w: rect.w, h: rect.h,
                        dp: dp, seq: seq++, z: z, i: parseInt(si),
                        cw: cr.width, ch: cr.height });
                }
            }
            return;
        }
        if (isFA(el)) {
            var fi = el.getAttribute('data-fi');
            if (fi !== null) out.icons.push({
                x: rect.x, y: rect.y, w: rect.w, h: rect.h,
                dp: dp, seq: seq++, z: z, i: parseInt(fi),
                cw: cr.width, ch: cr.height
            });
            return;
        }

        // Per-axis border-radius (CSS-spec: % is per-axis)
        var minDim = Math.min(rect.w, rect.h);
        var br_raw = cs.borderTopLeftRadius;
        var br_x = parseFloat(br_raw) || 0;
        var br_y = br_x;
        if (br_raw.indexOf('%') >= 0) {
            br_x = (br_x / 100) * rect.w;
            br_y = (br_y / 100) * rect.h;
        }
        var br_legacy = Math.min(br_x, br_y);  // back-compat single-radius
        // Circle classification: element must be roughly square AND radius ~50%.
        var aspect = Math.max(rect.w, rect.h) / Math.max(minDim, 0.001);
        var pct_x = rect.w > 0 ? (br_x / rect.w) : 0;
        var pct_y = rect.h > 0 ? (br_y / rect.h) : 0;
        var isCircle = (pct_x >= 0.45 && pct_y >= 0.45) && (aspect < 1.15) && (minDim > 20);

        var bg = cs.backgroundColor;
        var pushedBgShape = false;
        if (bg && bg !== 'rgba(0, 0, 0, 0)' && bg !== 'transparent') {
            var shapeEntry = {
                x: rect.x, y: rect.y, w: rect.w, h: rect.h,
                dp: dp, seq: seq++, z: z,
                bg: bg, op: parseFloat(cs.opacity),
                br: br_legacy, brX: br_x, brY: br_y
            };
            if (isCircle) shapeEntry.circ = true;
            out.shapes.push(shapeEntry);
            pushedBgShape = true;
        }

        if (isCircle) {
            var circBw = 0, circBc = null;
            var sides = ['Top','Right','Bottom','Left'];
            for (var sidx = 0; sidx < sides.length; sidx++) {
                var sd = sides[sidx];
                var sbw = parseFloat(cs['border' + sd + 'Width']);
                var sbs = cs['border' + sd + 'Style'];
                var sbc = cs['border' + sd + 'Color'];
                if (sbw >= 1 && sbs !== 'none' && sbc && sbc !== 'rgba(0, 0, 0, 0)') {
                    if (sbw > circBw) circBw = sbw;
                    if (!circBc) circBc = sbc;
                }
            }
            if (circBw >= 1 && circBc) {
                if (pushedBgShape) {
                    out.shapes[out.shapes.length - 1].bco = circBc;
                    out.shapes[out.shapes.length - 1].bw = circBw;
                } else {
                    out.shapes.push({
                        x: rect.x, y: rect.y, w: rect.w, h: rect.h,
                        dp: dp, seq: seq++, z: z,
                        bg: 'rgba(0, 0, 0, 0)', op: parseFloat(cs.opacity),
                        circ: true, bco: circBc, bw: circBw
                    });
                }
            }
        } else {
            var borderDefs = [
                ['Top',    function(r, bw) { return {x:r.x, y:r.y, w:r.w, h:bw}; }],
                ['Bottom', function(r, bw) { return {x:r.x, y:r.y+r.h-bw, w:r.w, h:bw}; }],
                ['Left',   function(r, bw) { return {x:r.x, y:r.y, w:bw, h:r.h}; }],
                ['Right',  function(r, bw) { return {x:r.x+r.w-bw, y:r.y, w:bw, h:r.h}; }]
            ];
            for (var bi = 0; bi < borderDefs.length; bi++) {
                var prop = borderDefs[bi][0];
                var mkR = borderDefs[bi][1];
                var bw = parseFloat(cs['border' + prop + 'Width']);
                var bs = cs['border' + prop + 'Style'];
                var bc = cs['border' + prop + 'Color'];
                if (bw >= 1 && bs !== 'none' && bc && bc !== 'rgba(0, 0, 0, 0)') {
                    var sr2 = mkR(rect, bw);
                    out.shapes.push({
                        x: sr2.x, y: sr2.y, w: sr2.w, h: sr2.h,
                        dp: dp + 0.5, seq: seq++, z: z,
                        bg: bc, op: parseFloat(cs.opacity), br: 0
                    });
                }
            }
        }

        // Text runs with precise bounding rects via Range API + multiline detection.
        var runs = [];
        var tMinX = Infinity, tMinY = Infinity, tMaxX = -Infinity, tMaxY = -Infinity;
        var hasTB = false;
        var lineCount = 0;
        var childNodes = el.childNodes;
        for (var ci2 = 0; ci2 < childNodes.length; ci2++) {
            var ch2 = childNodes[ci2];
            if (ch2.nodeType === 3) {
                var t = ch2.textContent.replace(/\s+/g, ' ').trim();
                if (t) {
                    runs.push(Object.assign({ t: t }, sty(el)));
                    var rng = document.createRange();
                    rng.selectNode(ch2);
                    var rr = rng.getBoundingClientRect();
                    if (rr.width > 0 && rr.height > 0) {
                        tMinX = Math.min(tMinX, rr.left);
                        tMinY = Math.min(tMinY, rr.top);
                        tMaxX = Math.max(tMaxX, rr.right);
                        tMaxY = Math.max(tMaxY, rr.bottom);
                        hasTB = true;
                        try {
                            var rects = rng.getClientRects();
                            if (rects && rects.length > lineCount) lineCount = rects.length;
                        } catch (e) {}
                    }
                }
            } else if (ch2.nodeType === 1 && inl(ch2) && vis(ch2) && !isFA(ch2) &&
                       ch2.tagName !== 'svg' && ch2.tagName !== 'SVG') {
                var inlBg = getComputedStyle(ch2).backgroundColor;
                var t2 = ch2.textContent.trim();
                if (t2) {
                    var rd = Object.assign({ t: t2 }, sty(ch2));
                    if (inlBg && inlBg !== 'rgba(0, 0, 0, 0)' && inlBg !== 'transparent') {
                        var inlRect = ch2.getBoundingClientRect();
                        rd.hlBg = inlBg;
                        rd.hlX = inlRect.left - ox;
                        rd.hlY = inlRect.top - oy;
                        rd.hlW = inlRect.width;
                        rd.hlH = inlRect.height;
                        rd.hlBr = parseInt(getComputedStyle(ch2).borderRadius) || 0;
                    }
                    runs.push(rd);
                    var rr2 = ch2.getBoundingClientRect();
                    if (rr2.width > 0 && rr2.height > 0) {
                        tMinX = Math.min(tMinX, rr2.left);
                        tMinY = Math.min(tMinY, rr2.top);
                        tMaxX = Math.max(tMaxX, rr2.right);
                        tMaxY = Math.max(tMaxY, rr2.bottom);
                        hasTB = true;
                    }
                }
            }
        }
        if (runs.length > 0) {
            var ty = hasTB ? (tMinY - oy) : rect.y;
            var th = hasTB ? (tMaxY - tMinY) : rect.h;
            var tx = hasTB ? (tMinX - ox) : rect.x;
            var tw = hasTB ? (tMaxX - tMinX) : rect.w;
            out.texts.push({
                x: rect.x, y: ty, w: rect.w, h: th,
                tx: tx, tw: tw,
                dp: dp, seq: seq++, z: z,
                runs: runs, ta: cs.textAlign,
                multiline: lineCount > 1
            });
        }

        dp++;
        var children = el.children;
        for (var ki = 0; ki < children.length; ki++) {
            var kid = children[ki];
            if (kid.nodeType !== 1) continue;
            if (isFA(kid) || kid.tagName === 'svg' || kid.tagName === 'SVG') {
                walk(kid);
                continue;
            }
            if (!inl(kid)) {
                walk(kid);
            } else {
                var hasBlock = false;
                for (var gi = 0; gi < kid.children.length; gi++) {
                    if (!inl(kid.children[gi])) { hasBlock = true; break; }
                }
                if (hasBlock) walk(kid);
            }
        }
        dp--;
    }

    walk(c);

    // Font width ratios: web font vs Windows fallback.
    var prevRatios = (window.__h2pFontRatios && typeof window.__h2pFontRatios === 'object') ? window.__h2pFontRatios : {};
    var _fonts = {};
    for (var _ti = 0; _ti < out.texts.length; _ti++) {
        for (var _ri = 0; _ri < out.texts[_ti].runs.length; _ri++) {
            var _ff = out.texts[_ti].runs[_ri].ff;
            if (_ff) _fonts[_ff] = 1;
        }
    }
    var _SYS = {'Segoe UI':1,'Arial':1,'Calibri':1,'Times New Roman':1,
                'Consolas':1,'Courier New':1,'Verdana':1,'Tahoma':1,'Georgia':1,'Trebuchet MS':1};
    var _REF = 'ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789';
    out.fontRatios = {};
    var _fnames = Object.keys(_fonts);
    for (var _fi = 0; _fi < _fnames.length; _fi++) {
        var _fn = _fnames[_fi];
        if (_SYS[_fn]) continue;
        if (prevRatios[_fn]) { out.fontRatios[_fn] = prevRatios[_fn]; continue; }
        var _s1 = document.createElement('span');
        _s1.style.cssText = 'position:absolute;visibility:hidden;white-space:nowrap;font-size:100px;font-family:"' + _fn + '",sans-serif';
        _s1.textContent = _REF;
        document.body.appendChild(_s1);
        var _mi = document.createElement('span');
        _mi.style.cssText = 'position:absolute;visibility:hidden;white-space:nowrap;font-size:100px;font-family:"' + _fn + '"';
        _mi.textContent = 'iiiiii';
        var _mw = document.createElement('span');
        _mw.style.cssText = 'position:absolute;visibility:hidden;white-space:nowrap;font-size:100px;font-family:"' + _fn + '"';
        _mw.textContent = 'MMMMMM';
        document.body.appendChild(_mi);
        document.body.appendChild(_mw);
        var _isMono = Math.abs(_mi.getBoundingClientRect().width - _mw.getBoundingClientRect().width) < 2;
        document.body.removeChild(_mi);
        document.body.removeChild(_mw);
        var _fb = _isMono ? 'Consolas' : 'Segoe UI';
        var _s2 = document.createElement('span');
        _s2.style.cssText = 'position:absolute;visibility:hidden;white-space:nowrap;font-size:100px;font-family:"' + _fb + '",sans-serif';
        _s2.textContent = _REF;
        document.body.appendChild(_s2);
        var _w1 = _s1.getBoundingClientRect().width;
        var _w2 = _s2.getBoundingClientRect().width;
        document.body.removeChild(_s1);
        document.body.removeChild(_s2);
        if (_w1 > 0 && _w2 > 0) {
            out.fontRatios[_fn] = { r: _w1 / _w2, fb: _fb };
        }
    }
    // Carry forward to next slide (cuts per-slide font-measurement DOM thrash).
    window.__h2pFontRatios = Object.assign({}, prevRatios, out.fontRatios);

    return out;
}"""


# ── Screenshot capture (BytesIO; no per-slide tempdir round-trip) ────

def _screenshot_with_isolation(page, selector: str) -> Optional[bytes]:
    """Hide ancestor siblings, screenshot the element, restore. Returns PNG bytes."""
    try:
        page.evaluate("(sel) => window.__h2pH(sel)", selector)
    except PlaywrightError as e:
        log.debug("hide failed for %s: %s", selector, type(e).__name__)
        return None
    try:
        handle = page.query_selector(selector)
        if not handle:
            return None
        return handle.screenshot()
    finally:
        try:
            page.evaluate("() => window.__h2pR()")
        except PlaywrightError:
            pass


def _screenshot_elements(page, data: WalkerOutput, ctx: SlideContext, report: ConversionReport, slide_name: str) -> list[dict]:
    """Capture per-element PNG bytes for SVGs and FA icons that fell back to screenshot."""
    images: list[dict] = []

    # Use the container width/height that the walker recorded (HIGH-6).
    cw = data.get("cw", ctx.vp_w)
    ch = data.get("ch", ctx.vp_h)

    for svg in data.get("svgs", []):
        # Full-slide decorative SVGs (>80% of CONTAINER, matching JS skip threshold).
        if svg["w"] > cw * SVG_FULL_SLIDE_THRESHOLD and svg["h"] > ch * SVG_FULL_SLIDE_THRESHOLD:
            continue
        si = int(svg["i"])
        try:
            png = _screenshot_with_isolation(page, f'[data-si="{si}"]')
            if png:
                images.append({**svg, "png": png})
            else:
                report.screenshots_failed += 1
                report.warn(slide_name, f"screenshot returned no bytes for svg {si}")
        except (PlaywrightError, OSError) as e:
            report.screenshots_failed += 1
            report.warn(slide_name, f"svg {si} screenshot failed: {type(e).__name__}")

    for icon in data.get("icons", []):
        fi = int(icon["i"])
        try:
            png = _screenshot_with_isolation(page, f'[data-fi="{fi}"]')
            if png:
                images.append({**icon, "png": png})
            else:
                report.screenshots_failed += 1
                report.warn(slide_name, f"screenshot returned no bytes for icon {fi}")
        except (PlaywrightError, OSError) as e:
            report.screenshots_failed += 1
            report.warn(slide_name, f"icon {fi} screenshot failed: {type(e).__name__}")

    return images


# ── Build PPTX slide from extracted data ─────────────────────────────

def _flatten_native_svgs(native_svgs: list[dict]) -> list[tuple[float, int, dict]]:
    """Return (dp, seq, prim) tuples for each SVG primitive so they can interleave
    with surrounding HTML shapes in the final z-order sort (HIGH-2 / LOG-C4)."""
    out: list[tuple[float, int, dict]] = []
    for nsvg in native_svgs:
        elements = nsvg.get("elements", [])
        for prim in elements:
            out.append((prim.get("dp", 0), prim.get("seq", 0), prim))
    return out


def build_slide(prs, data: WalkerOutput, images: list[dict], ctx: SlideContext, report: ConversionReport, slide_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    bg = parse_rgba(data.get("bg"))
    if not bg or bg[3] < MIN_CONTAINER_BG_ALPHA:
        bg = parse_rgba(data.get("bodyBg"))
    if bg:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = to_rgb(bg)

    # Merge all visual elements with a uniform sort key: (z, dp, seq).
    # CSS z-index now influences paint order. Native SVG primitives are flattened
    # so they can interleave with surrounding shapes.
    elems: list[tuple[str, tuple[int, float, int], dict]] = []
    for s in data.get("shapes", []):
        if s["w"] > ctx.vp_w - BG_DEDUP_MARGIN_PX and s["h"] > ctx.vp_h - BG_DEDUP_MARGIN_PX:
            sbg = parse_rgba(s["bg"])
            if sbg and bg and sbg[:3] == bg[:3]:
                continue
        key = (int(s.get("z", 0)), float(s.get("dp", 0)), int(s.get("seq", 0)))
        elems.append(("shape", key, s))
    for t in data.get("texts", []):
        key = (int(t.get("z", 0)), float(t.get("dp", 0)), int(t.get("seq", 0)))
        elems.append(("text", key, t))
    for i in images:
        key = (int(i.get("z", 0)), float(i.get("dp", 0)), int(i.get("seq", 0)))
        elems.append(("image", key, i))
    for prim_dp, prim_seq, prim in _flatten_native_svgs(data.get("nativeSvgs", [])):
        key = (0, float(prim_dp), int(prim_seq))
        elems.append(("nativesvg-prim", key, prim))

    # Clip to viewport.
    clipped: list[tuple[str, tuple[int, float, int], dict]] = []
    for etype, sort_key, ed in elems:
        ey = ed.get("y", 0)
        eh = ed.get("h", 0)
        ex = ed.get("x", 0)
        ew = ed.get("w", 0)
        if ey >= ctx.vp_h or ey + eh <= 0 or ex >= ctx.vp_w or ex + ew <= 0:
            continue
        if ex < 0:
            ed = {**ed, "w": ew + ex, "x": 0}
            ew = ed["w"]
            ex = 0
        if ey < 0:
            ed = {**ed, "h": eh + ey, "y": 0}
            eh = ed["h"]
            ey = 0
        if ey + eh > ctx.vp_h:
            ed = {**ed, "h": ctx.vp_h - ey}
            eh = ed["h"]
        if etype in ("shape", "image", "nativesvg-prim"):
            if ex + ew > ctx.vp_w:
                ed = {**ed, "w": ctx.vp_w - ex}
        elif etype == "text":
            visible_h = min(ctx.vp_h - ey, eh)
            if visible_h < eh * TEXT_OVERFLOW_DROP_RATIO:
                continue
        clipped.append((etype, sort_key, ed))

    clipped.sort(key=lambda e: e[1])

    font_ratios = data.get("fontRatios", {})

    for etype, _, ed in clipped:
        try:
            if etype == "shape":
                _add_shape(slide, ed, ctx)
            elif etype == "text":
                _add_text(slide, ed, ctx, font_ratios)
            elif etype == "image":
                _add_image(slide, ed, ctx)
            elif etype == "nativesvg-prim":
                _dispatch_native_svg_prim(slide, ed, ctx, font_ratios, report, slide_name)
        except (KeyError, TypeError, ValueError, AttributeError) as e:
            report.elements_skipped += 1
            report.warn(slide_name, f"skipped {etype}: {type(e).__name__}: {e}")


def _add_shape(slide, s: ShapeEntry, ctx: SlideContext) -> None:
    w, h = max(s["w"], 1), max(s["h"], 1)
    if w < MIN_SHAPE_PX and h < MIN_SHAPE_PX:
        return
    fill = parse_rgba(s["bg"])
    has_visible_fill = fill is not None and (s.get("op", 1.0) * fill[3]) >= MIN_VISIBLE_FILL_ALPHA
    border_color = parse_rgba(s.get("bco"))
    has_border = border_color is not None
    if not has_visible_fill and not has_border:
        return

    # Use OVAL only when the element is genuinely circular per CSS-spec percentages
    # (per-axis radius >= ~45% of width AND height, near-square aspect). Non-square
    # pills with high border-radius now correctly route to ROUNDED_RECTANGLE rather
    # than rendering as full ellipses.
    if s.get("circ", False):
        shape_type = MSO_SHAPE.OVAL
    elif s.get("br", 0) > ROUNDED_RECT_THRESHOLD_PX:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE

    shp = slide.shapes.add_shape(shape_type, ctx.px(s["x"]), ctx.px(s["y"]), ctx.px(w), ctx.px(h))
    if has_visible_fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = to_rgb(fill)
    else:
        shp.fill.background()
    if has_border:
        shp.line.color.rgb = to_rgb(border_color)
        shp.line.width = Pt(max(s.get("bw", 1) * PX_TO_PT, 0.25))
    else:
        shp.line.fill.background()


def _add_text(slide, t: TextEntry, ctx: SlideContext, font_ratios: dict) -> None:
    runs = t.get("runs", [])
    if not runs or t["w"] < MIN_TEXT_PX:
        return

    x, y, w, h = t["x"], t["y"], t["w"], t["h"]
    ta = t.get("ta", "left")
    max_fs = max(rd.get("fs", DEFAULT_FONT_SIZE_PX) for rd in runs)

    tx = t.get("tx", x)
    tw = t.get("tw", w)

    # Browser-precise multi-line detection (HIGH-4) -- falls back to legacy heuristic
    # if walker didn't emit it (older payload).
    is_multiline = bool(t.get("multiline", h > max_fs * 1.8))

    primary_font = runs[0].get("ff", "Segoe UI") if runs else "Segoe UI"
    _, fallback_factor = _resolve_font(primary_font, font_ratios)

    if is_multiline:
        w_use = w * fallback_factor
        if ta == "center":
            extra = w_use - w
            desired_x = x - extra / 2
            if desired_x < 0:
                w = max(1, w_use + desired_x)  # reduce w to preserve original right edge (HIGH-8)
                x = 0
            else:
                x = desired_x
                w = min(ctx.vp_w, w_use)
        else:
            x = t["x"]
            w = min(ctx.vp_w - x, w_use)
    else:
        if ta == "center":
            needed_w = w * fallback_factor
            extra = needed_w - w
            desired_x = x - extra / 2
            if desired_x < 0:
                w = max(1, needed_w + desired_x)
                x = 0
            else:
                x = desired_x
                w = min(ctx.vp_w, needed_w)
        else:
            x = tx
            w = min(ctx.vp_w - x, tw * fallback_factor)

    h = max(h, max_fs * MIN_TEXT_HEIGHT_FACTOR)
    if y + h > ctx.vp_h:
        h = max(ctx.vp_h - y, 1)

    txBox = slide.shapes.add_textbox(ctx.px(x), ctx.px(y), ctx.px(w), ctx.px(h))
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = bool(is_multiline)
    p = tf.paragraphs[0]
    p.alignment = ALIGN.get(ta, PP_ALIGN.LEFT)

    for rd in runs:
        if rd.get("hlBg"):
            hl_bg = parse_rgba(rd["hlBg"])
            if not hl_bg:
                continue
            hl_x, hl_y = rd["hlX"], rd["hlY"]
            hl_w, hl_h = rd["hlW"], rd["hlH"]
            hl_text = _strip_dangerous_unicode(_apply_text_transform(rd["t"], rd.get("tt", "none")))
            hlBox = slide.shapes.add_textbox(ctx.px(hl_x), ctx.px(hl_y), ctx.px(hl_w), ctx.px(hl_h))
            hlTf = hlBox.text_frame
            hlTf.auto_size = MSO_AUTO_SIZE.NONE
            hlTf.word_wrap = False
            hlTf.margin_left = hlTf.margin_right = hlTf.margin_top = hlTf.margin_bottom = 0
            hlBox.fill.solid()
            hlBox.fill.fore_color.rgb = to_rgb(hl_bg)
            hlP = hlTf.paragraphs[0]
            hlP.alignment = PP_ALIGN.CENTER
            hlR = hlP.add_run()
            hlR.text = hl_text
            font_name, _ = _resolve_font(rd.get("ff", "Segoe UI"), font_ratios)
            hlR.font.name = font_name
            hlR.font.size = ctx.pt(rd.get("fs", DEFAULT_FONT_SIZE_PX))
            hlR.font.bold = rd.get("fw", DEFAULT_FONT_WEIGHT) >= BOLD_WEIGHT_THRESHOLD
            co = parse_rgba(rd.get("co", ""))
            if co:
                hlR.font.color.rgb = to_rgb(co)

    for i, rd in enumerate(runs):
        text = _strip_dangerous_unicode(_apply_text_transform(rd["t"], rd.get("tt", "none")))
        if i > 0:
            text = " " + text
        r = p.add_run()
        r.text = text
        font_name, _ = _resolve_font(rd.get("ff", "Segoe UI"), font_ratios)
        r.font.name = font_name
        r.font.size = ctx.pt(rd.get("fs", DEFAULT_FONT_SIZE_PX))
        r.font.bold = rd.get("fw", DEFAULT_FONT_WEIGHT) >= BOLD_WEIGHT_THRESHOLD
        r.font.italic = rd.get("fi", False)
        co = parse_rgba(rd.get("co", ""))
        if co:
            if rd.get("hlBg"):
                hl_bg = parse_rgba(rd["hlBg"])
                if hl_bg:
                    r.font.color.rgb = to_rgb(hl_bg)  # invisible in main textbox; visible textbox renders separately
                else:
                    r.font.color.rgb = to_rgb(co)
            else:
                r.font.color.rgb = to_rgb(co)


def _add_image(slide, img: dict, ctx: SlideContext) -> None:
    png = img.get("png")
    if not png:
        return
    if img["w"] < MIN_IMG_PX or img["h"] < MIN_IMG_PX:
        return
    slide.shapes.add_picture(BytesIO(png), ctx.px(img["x"]), ctx.px(img["y"]), ctx.px(img["w"]), ctx.px(img["h"]))


# ── Native SVG primitives ────────────────────────────────────────────

def _apply_svg_style(shape, elem: SvgPrimitive) -> None:
    fill_color = parse_rgba(elem.get("fill"))
    stroke_color = parse_rgba(elem.get("stroke"))
    stroke_w = elem.get("strokeWidth", 0) or 0
    opacity = elem.get("opacity", 1.0)
    fill_op = elem.get("fillOpacity", 1.0)
    stroke_op = elem.get("strokeOpacity", 1.0)

    if fill_color and fill_color[3] * fill_op * opacity >= MIN_SVG_ALPHA:
        shape.fill.solid()
        shape.fill.fore_color.rgb = to_rgb(fill_color)
    else:
        shape.fill.background()

    if stroke_color and stroke_w > 0 and stroke_color[3] * stroke_op * opacity >= MIN_SVG_ALPHA:
        shape.line.color.rgb = to_rgb(stroke_color)
        shape.line.width = Pt(max(stroke_w * PX_TO_PT, 0.25))
        if elem.get("dashed"):
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    else:
        shape.line.fill.background()


def _add_svg_circle(slide, elem: SvgPrimitive, ctx: SlideContext) -> None:
    x, y, w, h = elem["x"], elem["y"], elem["w"], elem["h"]
    if w < MIN_SVG_RECT_PX and h < MIN_SVG_RECT_PX:
        return
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, ctx.px(x), ctx.px(y), ctx.px(max(w, 1)), ctx.px(max(h, 1)))
    _apply_svg_style(shp, elem)


def _add_svg_rect(slide, elem: SvgPrimitive, ctx: SlideContext) -> None:
    x, y, w, h = elem["x"], elem["y"], elem["w"], elem["h"]
    if w < MIN_SVG_RECT_PX and h < MIN_SVG_RECT_PX:
        return
    rx = elem.get("rx", 0)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > SVG_RECT_RX_THRESHOLD_PX else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, ctx.px(x), ctx.px(y), ctx.px(max(w, 1)), ctx.px(max(h, 1)))
    _apply_svg_style(shp, elem)


def _add_svg_freeform(slide, elem: SvgPrimitive, ctx: SlideContext) -> None:
    points = elem.get("points", [])
    if len(points) < 2:
        return
    closed = elem.get("closed", False)

    xs = [p[0] for p in points]
    ys = [p[1] for p in points]
    min_x, min_y = min(xs), min(ys)
    if max(xs) - min_x < MIN_FREEFORM_BBOX_PX and max(ys) - min_y < MIN_FREEFORM_BBOX_PX:
        return

    local = [(round(p[0] - min_x), round(p[1] - min_y)) for p in points]

    fb = slide.shapes.build_freeform(start_x=local[0][0], start_y=local[0][1], scale=ctx.scale)
    fb.add_line_segments(local[1:], close=closed)
    shape = fb.convert_to_shape(origin_x=ctx.px(min_x), origin_y=ctx.px(min_y))
    _apply_svg_style(shape, elem)


def _add_svg_text(slide, elem: SvgPrimitive, ctx: SlideContext, font_ratios: dict) -> None:
    text = (elem.get("text") or "").strip()
    if not text:
        return
    text = _strip_dangerous_unicode(text)
    x, y, w, h = elem["x"], elem["y"], elem["w"], elem["h"]
    ff = elem.get("fontFamily", "Segoe UI")
    font_name, fallback_factor = _resolve_font(ff, font_ratios)
    pad = SVG_TEXT_PADDING_PX
    w_use = w * fallback_factor
    txBox = slide.shapes.add_textbox(ctx.px(x - pad), ctx.px(y), ctx.px(w_use + pad * 2), ctx.px(max(h, 1)))
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    anchor = elem.get("textAnchor", "start")
    if anchor == "middle":
        p.alignment = PP_ALIGN.CENTER
    elif anchor == "end":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    r = p.add_run()
    r.text = text
    r.font.name = font_name
    fs = elem.get("fontSize", DEFAULT_FONT_SIZE_PX) or DEFAULT_FONT_SIZE_PX
    r.font.size = ctx.pt(fs)
    r.font.bold = elem.get("fontWeight", DEFAULT_FONT_WEIGHT) >= BOLD_WEIGHT_THRESHOLD
    fill_c = parse_rgba(elem.get("fill"))
    if fill_c:
        r.font.color.rgb = to_rgb(fill_c)


def _dispatch_native_svg_prim(slide, prim: SvgPrimitive, ctx: SlideContext, font_ratios: dict, report: ConversionReport, slide_name: str) -> None:
    etype = prim.get("type")
    try:
        if etype in ("circle", "ellipse"):
            _add_svg_circle(slide, prim, ctx)
        elif etype == "rect":
            _add_svg_rect(slide, prim, ctx)
        elif etype in ("line", "path", "polygon", "polyline"):
            _add_svg_freeform(slide, prim, ctx)
        elif etype == "text":
            _add_svg_text(slide, prim, ctx, font_ratios)
    except (KeyError, TypeError, ValueError, AttributeError) as e:
        report.elements_skipped += 1
        report.warn(slide_name, f"skipped SVG {etype}: {type(e).__name__}: {e}")


# ── HTML patching (input -> patched HTML for Chromium) ────────────────

def _patch_html(content: str, simplify_level: int) -> str:
    """Apply Tailwind v2 CDN swap + flex fix + optional CSS simplification."""
    patched = TAILWIND_V2_LINK_RE.sub(TAILWIND_CDN_TAG, content)
    inject = FLEX_FIX_CSS + _simplify_css(simplify_level)
    if "</head>" in patched:
        return patched.replace("</head>", inject + "</head>", 1)
    if "</HEAD>" in patched:
        return patched.replace("</HEAD>", inject + "</HEAD>", 1)
    if "<body" in patched:
        return patched.replace("<body", inject + "<body", 1)
    if "<BODY" in patched:
        return patched.replace("<BODY", inject + "<BODY", 1)
    return inject + patched


# ── Network-egress route handler ─────────────────────────────────────

def _make_route_handler(allowlist: list[str]):
    """Return a Playwright route callback that aborts non-allowlisted requests."""
    allowset = [a.lower() for a in allowlist]

    def handler(route, request):
        url = request.url.lower()
        if url.startswith("file://") or url.startswith("data:") or url.startswith("blob:"):
            route.continue_()
            return
        if any(host in url for host in allowset):
            route.continue_()
            return
        route.abort()

    return handler


# ── Slide ordering ───────────────────────────────────────────────────

_NATSORT_RE = re.compile(r"(\d+)")


def _natural_sort_key(stem: str) -> tuple:
    """Natural-order sort: 1 < 2 < 10, mixed alphanumerics handled sensibly."""
    parts = _NATSORT_RE.split(stem)
    return tuple(int(p) if p.isdigit() else p.lower() for p in parts)


# ── Conversion entry point (library-callable) ─────────────────────────

def convert(
    *,
    input_dir: Path,
    output: Path,
    width: int = DEFAULT_VP_W,
    height: int = DEFAULT_VP_H,
    simplify: int = 0,
    strict: bool = False,
    allow_javascript: bool = True,
    allow_network: Optional[list[str]] = None,
    tailwind_wait_ms: int = DEFAULT_TAILWIND_WAIT_MS,
    network_idle_timeout_ms: int = NETWORK_IDLE_TIMEOUT_MS,
    max_slides: int = MAX_SLIDES_DEFAULT,
) -> ConversionReport:
    """Convert HTML slides in ``input_dir`` to a PPTX at ``output``.

    Returns a ConversionReport with per-run statistics.
    Raises SystemExit on configuration errors.
    """
    ctx = SlideContext.from_viewport(width, height)
    report = ConversionReport()

    files = sorted(
        (f for f in input_dir.glob("*.html") if _safe_input_file(f, input_dir)),
        key=lambda f: _natural_sort_key(f.stem),
    )[:max_slides]
    report.slides_total = len(files)
    if not files:
        log.warning("No HTML files found in %s", input_dir)
        return report

    allowlist = list(allow_network) if allow_network is not None else list(DEFAULT_NETWORK_ALLOWLIST)
    log.info("Network allowlist: %s", ", ".join(allowlist) or "(empty)")

    prs = Presentation()
    prs.slide_width = ctx.slide_w
    prs.slide_height = ctx.slide_h

    with sync_playwright() as pw:
        browser = pw.chromium.launch(args=[
            "--disable-extensions",
            "--disable-background-networking",
            "--disable-features=ServiceWorker,SharedArrayBuffer",
            "--disable-default-apps",
            "--disable-sync",
            "--mute-audio",
        ], chromium_sandbox=True)
        try:
            for idx, hf in enumerate(files):
                slide_name = hf.name
                log.info("[%d/%d] %s", idx + 1, len(files), slide_name)
                # Per-slide isolated context: clears cookies/localStorage/SW between slides.
                context = browser.new_context(
                    viewport={"width": ctx.vp_w, "height": ctx.vp_h},
                    java_script_enabled=allow_javascript,
                )
                try:
                    context.route("**/*", _make_route_handler(allowlist))
                    context.add_init_script(HIDE_RESTORE_INIT_JS)
                    page = context.new_page()
                    try:
                        try:
                            html_content = hf.read_text(encoding="utf-8")
                        except (OSError, UnicodeDecodeError) as e:
                            report.slides_skipped += 1
                            report.warn(slide_name, f"read failed: {type(e).__name__}")
                            continue
                        if len(html_content.encode("utf-8", errors="ignore")) > MAX_INPUT_HTML_BYTES:
                            report.slides_skipped += 1
                            report.warn(slide_name, f"file exceeds {MAX_INPUT_HTML_BYTES} bytes")
                            continue
                        patched = _patch_html(html_content, simplify)
                        # set_content avoids a tempdir round-trip vs goto(file://).
                        try:
                            page.set_content(patched, wait_until="domcontentloaded", timeout=15000)
                        except PlaywrightTimeoutError:
                            report.slides_skipped += 1
                            report.warn(slide_name, "set_content timeout")
                            continue
                        try:
                            page.wait_for_load_state("networkidle", timeout=network_idle_timeout_ms)
                        except PlaywrightTimeoutError:
                            report.networkidle_timeouts += 1
                        # Tailwind JIT readiness probe: replaces the legacy unconditional 500 ms wait.
                        try:
                            page.wait_for_function(TAILWIND_READY_JS, timeout=tailwind_wait_ms)
                        except PlaywrightTimeoutError:
                            report.tailwind_wait_timeouts += 1
                        # Brief minimum wait for layout to settle (much shorter than legacy 500ms).
                        page.wait_for_timeout(DEFAULT_TAILWIND_MIN_MS)

                        try:
                            page.evaluate(PREPROCESS_JS)
                        except PlaywrightError as e:
                            report.preprocess_failures += 1
                            report.warn(slide_name, f"preprocess: {type(e).__name__}")

                        try:
                            data = page.evaluate(EXTRACT_JS)
                        except PlaywrightError as e:
                            report.slides_skipped += 1
                            report.warn(slide_name, f"extract: {type(e).__name__}")
                            continue
                        if not data:
                            report.slides_skipped += 1
                            log.info("  skipped (no container)")
                            continue

                        images = _screenshot_elements(page, data, ctx, report, slide_name)
                        build_slide(prs, data, images, ctx, report, slide_name)
                        report.slides_processed += 1

                        ns = len(data.get("shapes", []))
                        nt = len(data.get("texts", []))
                        ni = len(images)
                        nn = sum(len(nsv.get("elements", [])) for nsv in data.get("nativeSvgs", []))
                        log.info("  %d shapes, %d texts, %d imgs%s", ns, nt, ni,
                                 f", {nn} svg" if nn else "")
                    finally:
                        page.close()
                finally:
                    context.close()
        finally:
            browser.close()

    # Save with atomic-rename via temp file, then fall back to .partial on failure.
    try:
        tmp = output.with_suffix(output.suffix + ".tmp")
        prs.save(str(tmp))
        try:
            os.replace(tmp, output)
        except OSError:
            # If atomic rename fails (Windows: target locked), fall back to direct save.
            prs.save(str(output))
    except (OSError, PermissionError) as e:
        log.error("save failed for %s: %s", output, type(e).__name__)
        fallback = output.with_suffix(".partial.pptx")
        try:
            prs.save(str(fallback))
            report.save_partial = True
            log.warning("Partial output saved: %s", fallback.name)
        except (OSError, PermissionError) as e2:
            log.error("unable to save any output: %s; fallback: %s", type(e).__name__, type(e2).__name__)
            raise SystemExit(1) from e2

    if strict and report.has_failures():
        log.warning("strict mode: failures detected -- exiting with code 2")
        raise SystemExit(2)

    return report


# ── CLI ─────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Convert HTML slides to native PPTX elements")
    parser.add_argument("-i", "--input", default=str(DEFAULT_HTML_DIR),
                        help="Directory with HTML files (default: presentazione_html)")
    parser.add_argument("-o", "--output", default=str(DEFAULT_OUTPUT),
                        help="Output PPTX path (default: Slides1.pptx)")
    parser.add_argument("--width", type=int, default=DEFAULT_VP_W,
                        help="Viewport width in pixels (default: %(default)s)")
    parser.add_argument("--height", type=int, default=DEFAULT_VP_H,
                        help="Viewport height in pixels (default: %(default)s)")
    parser.add_argument("-s", "--simplify", type=int, default=0, choices=range(0, 11),
                        metavar="0-10",
                        help="CSS simplification level (default: 0). "
                             "0=original, 1-2=no animations, 3-4=+no shadows, "
                             "5-6=+no filters, 7-8=+no bg-images/gradients, "
                             "9-10=+no outlines/pseudo-elements")
    parser.add_argument("--strict", action="store_true",
                        help="Exit with code 2 if any slide hits a warning or partial save.")
    parser.add_argument("--no-javascript", action="store_true",
                        help="Disable JavaScript in the headless browser (safer for untrusted HTML).")
    parser.add_argument("--allow-network", action="append", default=None, metavar="HOST",
                        help="Hostname substring to permit during conversion (repeatable). "
                             f"Default allowlist: {', '.join(DEFAULT_NETWORK_ALLOWLIST)}")
    parser.add_argument("--block-network", action="store_true",
                        help="Block ALL network requests (overrides --allow-network); "
                             "useful for offline / vendored decks.")
    parser.add_argument("--tailwind-wait-ms", type=int, default=DEFAULT_TAILWIND_WAIT_MS,
                        help="Max wait for Tailwind JIT readiness signal (ms; default %(default)s).")
    parser.add_argument("--max-slides", type=int, default=MAX_SLIDES_DEFAULT,
                        help="Hard cap on number of slides processed (default %(default)s).")
    parser.add_argument("-v", "--verbose", action="count", default=0,
                        help="Increase log verbosity (-v=info, -vv=debug).")
    return parser


def _configure_logging(verbosity: int) -> None:
    level = logging.WARNING
    if verbosity == 1:
        level = logging.INFO
    elif verbosity >= 2:
        level = logging.DEBUG
    logging.basicConfig(level=level, format="%(message)s", stream=sys.stderr)
    # Surface INFO to stderr but keep our progress lines on stdout via log.info().
    log.setLevel(level)


def main() -> None:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")  # type: ignore[attr-defined]

    parser = _build_parser()
    args = parser.parse_args()
    _configure_logging(args.verbose)

    if args.width <= 0 or args.height <= 0:
        parser.error("--width and --height must be positive integers")

    input_dir = _safe_input_dir(args.input)
    output = _safe_output_path(args.output)

    if args.block_network:
        allowlist: list[str] = []
    elif args.allow_network is not None:
        allowlist = args.allow_network
    else:
        allowlist = list(DEFAULT_NETWORK_ALLOWLIST)

    simp = args.simplify
    simp_label = {0: "off", 1: "light", 2: "light", 3: "medium", 4: "medium",
                  5: "medium-high", 6: "medium-high", 7: "heavy", 8: "heavy",
                  9: "maximum", 10: "maximum"}.get(simp, f"level {simp}")
    if simp > 0:
        log.info("CSS simplification: %d/10 (%s)", simp, simp_label)

    try:
        report = convert(
            input_dir=input_dir,
            output=output,
            width=args.width,
            height=args.height,
            simplify=simp,
            strict=args.strict,
            allow_javascript=not args.no_javascript,
            allow_network=allowlist,
            tailwind_wait_ms=args.tailwind_wait_ms,
            max_slides=args.max_slides,
        )
    except SystemExit:
        raise
    except (OSError, PermissionError) as e:
        log.error("%s: %s", type(e).__name__, e)
        sys.exit(1)

    if output.exists():
        kb = output.stat().st_size // 1024
        log.warning("Saved: %s (%s; %d KB)", output, report.summary_line(), kb)
    else:
        log.warning("Run finished: %s", report.summary_line())


if __name__ == "__main__":
    main()
